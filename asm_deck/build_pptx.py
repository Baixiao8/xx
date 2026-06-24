# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

PURPLE=RGBColor(0x5B,0x3E,0x70); PURPLE_DEEP=RGBColor(0x2E,0x1F,0x3D); PURPLE_MID=RGBColor(0x46,0x31,0x5C)
GOLD=RGBColor(0xE8,0xA9,0x3B); GOLD_SOFT=RGBColor(0xF2,0xC6,0x6B); LILAC=RGBColor(0xB9,0xA4,0xC9)
IVORY=RGBColor(0xF7,0xF3,0xEF); INK=RGBColor(0x2C,0x24,0x33); MUTED=RGBColor(0x6E,0x64,0x78)
LINE=RGBColor(0xE3,0xD9,0xEA); CREAM=RGBColor(0xFA,0xF6,0xF2); GOLDINK=RGBColor(0xA9,0x71,0x1A)
WHITE=RGBColor(0xFF,0xFF,0xFF); CREAMTX=RGBColor(0xEA,0xDF,0xF2); LILTX=RGBColor(0xCB,0xB9,0xDC)
FONT="Microsoft YaHei"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK=prs.slide_layouts[6]
A="assets/"

def sf(run,size,bold,color,name=FONT):
    run.font.size=Pt(size); run.font.bold=bold; run.font.color.rgb=color; run.font.name=name
    rPr=run._r.get_or_add_rPr()
    for tag in ('a:ea','a:cs'):
        e=rPr.find(qn(tag))
        if e is None:
            e=rPr.makeelement(qn(tag),{}); rPr.append(e)
        e.set('typeface',name)

def slide(bg=IVORY, scene=None):
    s=prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb=bg
    if scene:
        s.shapes.add_picture(A+scene+'.png',0,0,width=SW,height=SH)
    return s

def box(s,l,t,w,h):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tb,tf

def para(tf,first,segs,size=15,color=INK,bold=False,align=PP_ALIGN.LEFT,after=4,lh=1.25):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(after); p.space_before=Pt(0); p.line_spacing=lh
    if isinstance(segs,str): segs=[(segs,color,bold)]
    for seg in segs:
        txt=seg[0]; col=seg[1] if len(seg)>1 else color; bd=seg[2] if len(seg)>2 else bold
        sz=seg[3] if len(seg)>3 else size
        r=p.add_run(); r.text=txt; sf(r,sz,bd,col)
    return p

def kicker(s,text,color=GOLD,top=0.46):
    _,tf=box(s,0.62,top,11.8,0.4)
    para(tf,True,[("◆  ",GOLD,True,12),(text,color,True,13)],after=0)

def title(s,text,top=0.9,color=PURPLE,size=30):
    _,tf=box(s,0.62,top,12.0,1.1); para(tf,True,text,size=size,color=color,bold=True,after=0,lh=1.1)

def rule(s,top,left=0.64,w=0.78):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(left),Inches(top),Inches(w),Inches(0.055))
    r.fill.solid(); r.fill.fore_color.rgb=GOLD; r.line.fill.background(); r.shadow.inherit=False

def foot(s,num,dark=False):
    _,tf=box(s,0.5,7.05,3,0.35); para(tf,True,"芯火计划",size=10,color=LILAC if dark else LILAC,after=0)
    _,tf2=box(s,11.8,7.05,1.0,0.35); para(tf2,True,str(num),size=10,color=LILAC,align=PP_ALIGN.RIGHT,after=0)

def card(s,l,t,w,h,fill=WHITE,ln=LINE,accent=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.adjustments[0]=0.06
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    sh.line.color.rgb=ln; sh.line.width=Pt(0.75); sh.shadow.inherit=False
    if accent:
        ac=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l+0.0),Inches(t),Inches(w),Inches(0.08))
        ac.fill.solid(); ac.fill.fore_color.rgb=accent; ac.line.fill.background(); ac.shadow.inherit=False
    return sh

def tag(s,l,t,text,care=False):
    w=0.18+0.135*len(text)
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(0.32))
    sh.adjustments[0]=0.5
    sh.fill.solid(); sh.fill.fore_color.rgb= RGBColor(0xEF,0xE7,0xF1) if care else RGBColor(0xFB,0xF0,0xDC)
    sh.line.color.rgb= PURPLE if care else GOLD; sh.line.width=Pt(0.75); sh.shadow.inherit=False
    tf=sh.text_frame; tf.margin_top=0; tf.margin_bottom=0; tf.word_wrap=False
    para(tf,True,text,size=10.5,color= PURPLE if care else GOLDINK,bold=True,align=PP_ALIGN.CENTER,after=0)
    return l+w+0.12

def bullets(s,l,t,w,items,size=14.5,color=INK,gap=5,lh=1.22,h=4.0):
    _,tf=box(s,l,t,w,h)
    for i,it in enumerate(items):
        segs=[("◆  ",GOLD,True,size-2)]
        if isinstance(it,list): segs+=it
        else: segs.append((it,color,False,size))
        para(tf,i==0,segs,after=gap,lh=lh)
    return tf

def table(s,l,t,w,data,colw,fs=11,headfs=11.5,rh=None):
    rows=len(data); cols=len(data[0])
    gt=s.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(w),Inches(0.4*rows)).table
    gt.first_row=False; gt.horz_banding=False
    tot=sum(colw)
    for i in range(cols): gt.columns[i].width=Inches(w*colw[i]/tot)
    for r in range(rows):
        if rh: gt.rows[r].height=Inches(rh)
        for c in range(cols):
            cell=gt.cell(r,c)
            cell.margin_left=Inches(0.09); cell.margin_right=Inches(0.06)
            cell.margin_top=Inches(0.025); cell.margin_bottom=Inches(0.025)
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r==0:
                cell.fill.fore_color.rgb=PURPLE; col=WHITE; bd=True; sz=headfs
            else:
                cell.fill.fore_color.rgb= CREAM if r%2==1 else WHITE; col=INK; bd=False; sz=fs
            tf=cell.text_frame; tf.word_wrap=True
            p=tf.paragraphs[0]; p.line_spacing=1.05
            r1=p.add_run(); r1.text=data[r][c]; sf(r1,sz,bd,col)
    return gt

def pic_panel(s,l,t,w,h,img,panel=PURPLE_DEEP):
    p=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    p.adjustments[0]=0.05; p.fill.solid(); p.fill.fore_color.rgb=panel
    p.line.fill.background(); p.shadow.inherit=False
    from PIL import Image
    iw,ih=Image.open(A+img+'.png').size
    ar=iw/ih; pad=0.12; aw=w-2*pad; ah=h-2*pad
    if aw/ah>ar: dw=ah*ar; dh=ah
    else: dw=aw; dh=aw/ar
    s.shapes.add_picture(A+img+'.png',Inches(l+(w-dw)/2),Inches(t+(h-dh)/2),Inches(dw),Inches(dh))

# ============================== SLIDES ==============================
# 1 COVER
s=slide(PURPLE_DEEP,"scene-night")
_,tf=box(s,0.9,1.7,9,0.5); para(tf,True,[("◆  ",GOLD,True,13),("ASM × 木贡小学 · 2026 助学公益",GOLD,True,15)],after=0)
_,tf=box(s,0.85,2.25,11,1.5); para(tf,True,"芯火计划",size=66,color=WHITE,bold=True,after=0)
rule(s,3.95,0.9,1.0)
_,tf=box(s,0.9,4.25,11,1.4)
para(tf,True,"每个孩子，都是一束等待被点亮的光",size=24,color=GOLD_SOFT,bold=True,after=6)
para(tf,False,"— 一束芯火，照亮一座山 —",size=20,color=WHITE,bold=True,after=10)
para(tf,False,"贵州 · 六盘水 · 木贡小学　|　线下公益面对面专场　|　执行方案",size=14,color=LILTX)
foot(s,"",True)

# 2 项目概览
s=slide(); kicker(s,"项目概览　Project Overview"); title(s,"一次把“陪伴”做实的线下公益专场"); rule(s,1.78)
stats=[("3 天","重点设计 Day1 晚间破冰 + Day2 全天面对面公益"),
       ("20 + 20","ASM 工程师约 20 人（含女性工程师）× 木贡学生约 20 人，1:1 结对"),
       ("5 年","自 2021 定点帮扶至今，2026 升级为线下专场")]
for i,(big,desc) in enumerate(stats):
    l=0.64+i*4.05; card(s,l,2.0,3.8,1.55,accent=GOLD)
    _,tf=box(s,l+0.25,2.2,3.3,1.2); para(tf,True,big,size=30,color=GOLD,bold=True,after=4)
    para(tf,False,desc,size=12.5,color=INK,lh=1.25)
_,tf=box(s,0.64,3.9,6.0,0.4); para(tf,True,"背景",size=18,color=PURPLE,bold=True,after=2)
bullets(s,0.64,4.35,6.0,["木贡留守儿童占比高，科创、文体、线下陪伴资源稀缺",
    "以往以物资捐赠、远程研学为主；2026 落地校园线下专场",
    "融合科创科普 + 文体互动 + 物资捐赠，深化校企长效合作"],size=13.5,gap=4)
_,tf=box(s,6.95,3.9,6.0,0.4); para(tf,True,"本方案目标（可衡量）",size=18,color=PURPLE,bold=True,after=2)
bullets(s,6.95,4.35,6.0,["20 个孩子，每人亲手点亮一颗 LED、完成一件科创作品",
    "每个孩子获 1 位工程师全程 1:1 陪伴与一封“未来之信”",
    "留下：星河长卷 + 时间胶囊 + 运动角，延续到下次回访"],size=13.5,gap=4)
foot(s,2)

# 3 使命
s=slide(PURPLE_DEEP,"scene-radiant"); kicker(s,"为什么是 ASM 来做　第一性原理")
_,tf=box(s,0.62,1.5,7.6,2.0)
para(tf,True,[("「科技本身不改变世界，",WHITE,True,30)],after=2,lh=1.2)
para(tf,False,[("是人，改变世界。」",GOLD_SOFT,True,30)],after=6,lh=1.2)
para(tf,False,"—— ASM 官方使命：People are at the heart of ASM",size=14,color=LILTX)
bullets(s,0.62,3.9,7.7,[
    "ASM 造世界上最精密的芯片设备，但相信真正改变世界的是人",
    [("对“人”的投资，才是 ASM 使命最真实的表达",CREAMTX,False,15)],
    [("山区的孩子，正是",CREAMTX,False,15),("未来可能改变世界的人",GOLD_SOFT,True,15)],
    "→ 这不是顺便做善事，而是 ASM 使命最纯粹的一次落地"],size=15,color=CREAMTX,gap=7)
foot(s,3,True)

# 4 价值观DNA
s=slide(); kicker(s,"价值观地基　全程对齐，绝不偏离"); title(s,"每个环节，都长在 ASM 的价值观上"); rule(s,1.78)
vals=[("We Innovate · 创新","STEM 科研讲座、亲手点亮 LED、动手实验——激发好奇心与创新思维。"),
      ("We Deliver · 兑现","2021 至今 5 年定点帮扶；时间胶囊、长期之约——说到做到，持续兑现。"),
      ("We Care · 关怀","1:1 面对面陪伴、平等不俯视、孩子是主人——把陪伴做实。")]
for i,(h,b) in enumerate(vals):
    l=0.64+i*4.05; card(s,l,2.0,3.8,1.9,accent=GOLD)
    _,tf=box(s,l+0.25,2.25,3.3,1.5); para(tf,True,h,size=17,color=PURPLE,bold=True,after=6)
    para(tf,False,b,size=13,color=INK,lh=1.3)
card(s,0.64,4.2,12.05,1.5,fill=CREAM)
_,tf=box(s,0.9,4.42,11.5,1.1)
para(tf,True,[("WIN　",PURPLE,True,14),("Women's Initiative Network：「当有抱负、多元的头脑汇聚，才有开创性的进步」——本次由女性工程师把这份信念，从 ASM 员工延伸给大山里的女孩。",INK,False,13)],after=6,lh=1.3)
para(tf,False,[("STEM　",GOLDINK,True,14),("ASM 全球长期支持 STEM 与社区项目，与学校、NGO 合作——本次正是该战略的一环。",INK,False,13)],lh=1.3)
foot(s,4)

# 5 主题推导
s=slide(PURPLE_DEEP,"scene-lattice"); kicker(s,"主题创意推导　Theme Derivation")
title(s,"为什么叫「芯火计划」",color=WHITE); rule(s,1.78)
card(s,0.64,2.0,3.95,1.7,fill=PURPLE_MID,ln=PURPLE_MID)
_,tf=box(s,0.85,2.18,3.55,1.4); para(tf,True,"「芯」一字三义",size=16,color=GOLD_SOFT,bold=True,after=4)
para(tf,False,"芯＝芯片（ASM 基因）｜谐音 心（We Care 的温度）｜谐音 星，「芯火」即「星火」——可以燎原。",size=12.5,color=CREAMTX,lh=1.3)
card(s,4.75,2.0,3.95,1.7,fill=PURPLE_MID,ln=PURPLE_MID)
_,tf=box(s,4.96,2.18,3.55,1.4); para(tf,True,"方法＝ALD 逐层点亮",size=16,color=GOLD_SOFT,bold=True,after=4)
para(tf,False,"ASM 本业是 ALD 原子层沉积：伟大不是一蹴而就，而是一层一层、精准而耐心地生长出来。",size=12.5,color=CREAMTX,lh=1.3)
bullets(s,0.64,4.1,8.2,[
    [("每个孩子，都是一束等待被点亮的光",WHITE,True,15),("（前提：光本就在）",CREAMTX,False,14)],
    [("我们像做 ALD 一样点亮它",WHITE,True,15),("——一层陪伴、一层赋能，不急于求成",CREAMTX,False,14)],
    [("每束光连起来，照亮一座山",GOLD_SOFT,True,15),("（人汇聚成改变世界的力量）",CREAMTX,False,14)]],gap=8)
foot(s,5,True)

# 6 视觉母题
s=slide(); kicker(s,"视觉母题　Visual Key"); title(s,"让 ASM 的品牌符号，长进公益的视觉里"); rule(s,1.78)
hexs=s.shapes.add_shape(MSO_SHAPE.HEXAGON,Inches(1.4),Inches(2.5),Inches(2.6),Inches(2.6))
hexs.rotation=90; hexs.fill.background(); hexs.line.color.rgb=PURPLE; hexs.line.width=Pt(2.5); hexs.shadow.inherit=False
st=s.shapes.add_shape(MSO_SHAPE.STAR_6_POINT,Inches(1.95),Inches(3.05),Inches(1.5),Inches(1.5))
st.fill.solid(); st.fill.fore_color.rgb=GOLD; st.line.fill.background(); st.shadow.inherit=False
_,tf=box(s,1.0,5.3,3.4,0.5); para(tf,True,"ASM logo ＝ 六边形里的星形 / 晶格",size=13,color=MUTED,align=PP_ALIGN.CENTER)
bullets(s,5.0,2.45,7.6,[
    [("它像一颗 星",PURPLE,True,15),(" → 呼应“星河长卷”与”照亮一座山”",INK,False,15)],
    [("它又像 晶格 / 原子层",PURPLE,True,15),(" → 对上 ALD 逐层生长的技术内核",INK,False,15)],
    "孩子画的每一颗“星“ = logo 的星形几何（统一图形语言）",
    "星河上”星与星的连线” = logo 内部的线条结构",
    "每页以星形做隐性栅格与装饰，品牌符号自然生长，而非贴角标"],size=15,gap=7)
card(s,5.0,5.25,7.6,1.0,fill=CREAM)
_,tf=box(s,5.25,5.42,7.1,0.7); para(tf,True,[("三位一体：",PURPLE,True,14),("视觉（星河/晶格）＝ 技术（ALD）＝ 价值观（点亮人）。一眼看懂，无法复制。",INK,False,13.5)],lh=1.3)
foot(s,6)

# 7 情感主线 timeline
s=slide(PURPLE_DEEP); title(s,"从一束火，到照亮一座山",color=WHITE,top=0.9); kicker(s,"全景故事线　Emotional Arc"); rule(s,1.78)
steps=[("遇","芯·遇","Day1 晚·破冰·相遇"),("启","芯·启","Day2上午·STEM·开眼界"),
       ("她","芯·她","Day2午后·WIN·榜样"),("燃","芯·燃","Day2下午·运动·并肩"),("山","芯·山","收尾·升华·回响")]
n=len(steps); x0=0.95; gap=(12.0-0)/n
ly=Inches(2.95)
lnsh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1.45),Inches(2.92),Inches(10.4),Inches(0.04))
lnsh.fill.solid(); lnsh.fill.fore_color.rgb=GOLD; lnsh.line.fill.background(); lnsh.shadow.inherit=False
for i,(c,t,d) in enumerate(steps):
    cx=1.45+i*2.5
    o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx-0.42),Inches(2.5),Inches(0.84),Inches(0.84))
    o.fill.solid(); o.fill.fore_color.rgb=PURPLE_MID; o.line.color.rgb=GOLD; o.line.width=Pt(2.5); o.shadow.inherit=False
    tf=o.text_frame; para(tf,True,c,size=20,color=WHITE,bold=True,align=PP_ALIGN.CENTER,after=0)
    _,tf=box(s,cx-1.1,3.5,2.2,1.0)
    para(tf,True,t,size=16,color=WHITE,bold=True,align=PP_ALIGN.CENTER,after=2)
    para(tf,False,d,size=11.5,color=LILTX,align=PP_ALIGN.CENTER,lh=1.2)
card(s,0.64,4.75,5.9,1.5,fill=PURPLE_MID,ln=PURPLE_MID)
_,tf=box(s,0.9,4.95,5.4,1.15); para(tf,True,"五章以「芯·X」串联，每章让 slogan 回响一次；情绪层层升温，最终收束到“一束芯火，照亮一座山”。",size=14,color=CREAMTX,lh=1.35)
card(s,6.75,4.75,5.9,1.5,fill=PURPLE_MID,ln=PURPLE_MID)
_,tf=box(s,7.0,4.95,5.4,1.15); para(tf,True,"孩子内心：陌生 → “我有搭档” → “我也是工程师” → “女孩也能当领导” → “我能赢” → “我留下了光，未来会回响“。",size=14,color=CREAMTX,lh=1.35)
foot(s,7,True)

# 8 三天行程
s=slide(); kicker(s,"整体行程　3-Day Itinerary"); title(s,"三天总览（重点：Day1 晚 + Day2 全天）"); rule(s,1.78)
data=[["日期","时间","活动","地点"],
 ["Day1","17:00–18:00","抵达六盘水 · 入住休整","市区酒店"],
 ["Day1","18:00–19:30","芯·遇：公益破冰 + 结对晚宴（重点）","酒店会议室"],
 ["Day2","08:00–10:00","前往木贡小学（车程约 90 分钟）","大巴"],
 ["Day2","10:00–16:00","公益面对面：芯·启 / 芯·她 / 芯·燃 / 芯·山（重点）","木贡小学"],
 ["Day2","16:00–19:00","返回市区 · 团队晚餐复盘","市区"],
 ["Day3","09:00–17:00","返程（贵阳龙洞堡 → 上海浦东）","返程"]]
table(s,0.64,2.05,12.05,data,[1,1.6,4.2,1.6],fs=12.5,headfs=13,rh=0.5)
_,tf=box(s,0.64,6.2,12,0.5); para(tf,True,"本方案聚焦两个“重点时段”的玩法与执行；其余为行程衔接，按既定安排执行。",size=12,color=MUTED)
foot(s,8)

# 9 Day1 芯·遇
s=slide(); kicker(s,"DAY 1 · 18:00–19:30 · 酒店会议室"); title(s,"芯·遇 ｜ 相遇，是点亮的开始")
nx=tag(s,0.64,1.62,"We Care"); tag(s,nx,1.62,"人遇见人，光才亮",care=True); rule(s,2.1)
data=[["时间","环节 & 玩法","负责人","物料"],
 ["18:00–18:20","寻找点亮你的那个人：每人拿芯片造型徽章的一半，正确两半拼合才亮灯；找到搭档=三天”芯火搭档”","主持+全员","发光拼合徽章×40"],
 ["18:20–19:10","搭档共进晚餐 + 暖场小游戏（三个关键词 / 猜职业），快速破冰","主持","桌签、游戏卡"],
 ["19:10–19:30","点亮第一束芯火：女性“首席点灯人“点亮”芯火树”，留钩子——“明天你会亲手点亮自己的光”","首席点灯人","芯火树灯组"]]
table(s,0.64,2.35,12.05,data,[1.3,5.2,1.4,1.4],fs=12,headfs=12.5,rh=0.7)
_,tf=box(s,0.64,5.7,12,1.0); para(tf,True,[("话术示例：",PURPLE,True,13),("“今晚我们点亮了第一束火。明天，每个人都会亲手点亮属于自己的光；所有人的光连在一起，会照亮整座山。”",INK,False,13)],lh=1.35)
foot(s,9)

# 10 Day2 总览
s=slide(PURPLE_DEEP); kicker(s,"DAY 2 · 木贡小学 · 10:00–16:00"); title(s,"公益面对面 · 全天 Run-of-show",color=WHITE); rule(s,1.78)
data=[["时间","章节 / 环节","场地"],
 ["10:00–10:20","抵达 · 孩子带路（孩子当东道主，带工程师参观学校）","校园"],
 ["10:20–11:00","开幕 + 公益捐赠仪式（交火种 + 孩子回赠）","教室 / 礼堂"],
 ["11:00–12:00","芯·启：STEM 科研讲座 + 亲手点亮 LED","教室"],
 ["12:00–13:00","搭档共进午餐","食堂"],
 ["13:00–13:40","芯·她：WIN 女性领导力","教室"],
 ["13:40–15:10","芯·燃：趣味运动会 + 篮球友谊赛","操场 / 篮球场"],
 ["15:10–16:00","芯·山：共绘星河 + 传光仪式 + 时间胶囊","操场 / 教室"]]
table(s,0.64,2.05,12.05,data,[1.4,5.5,1.6],fs=12.5,headfs=13,rh=0.46)
foot(s,10,True)

# 11 芯·启
s=slide(); kicker(s,"DAY 2 · 10:20–12:00"); title(s,"芯·启 ｜ 亲手点亮你的第一束光")
nx=tag(s,0.64,1.62,"We Deliver"); nx=tag(s,nx,1.62,"We Innovate"); tag(s,nx,1.62,"STEM",care=True); rule(s,2.1)
card(s,0.64,2.35,5.9,2.7)
_,tf=box(s,0.9,2.55,5.4,2.4); para(tf,True,"① 开幕 + 公益捐赠仪式",size=17,color=PURPLE,bold=True,after=6)
bullets(s,0.9,3.0,5.4,["主语从”捐给你们“→”把火种交到你们手里”","助学金 + 科教物资 + STEM 套件捐赠",
    "孩子回赠画 / 信 → 孩子也成为给予者","女性点灯人 + 校方共同启动【对齐 We Deliver】"],size=13,gap=4)
card(s,6.75,2.35,5.9,2.7)
_,tf=box(s,7.0,2.55,5.4,2.4); para(tf,True,"② STEM 科研讲座 + 点亮 LED",size=17,color=PURPLE,bold=True,after=6)
bullets(s,7.0,3.0,5.4,["把 ASM 内部 STEM 课改成“山区孩子听得懂的科研课”","用”一层一层造芯片(ALD)”的通俗讲法",
    "1:1 带教，孩子亲手点亮一颗 LED","作品归孩子所有，可带回家展示【对齐 We Innovate】"],size=13,gap=4)
_,tf=box(s,0.64,5.35,12,1.0); para(tf,True,[("工程师话术：",PURPLE,True,13),("“你刚做的，和我们工厂里做的是同一件事——让芯片发光。这，就是你点亮的第一束芯火。你今天就是一名工程师。”",INK,False,13)],lh=1.35)
foot(s,11)

# 12 芯·她
s=slide(PURPLE_DEEP,"scene-radiant"); kicker(s,"DAY 2 · 13:00–13:40"); title(s,"芯·她 ｜ 她的故事，你的可能",color=WHITE)
nx=tag(s,0.64,1.62,"We Care"); tag(s,nx,1.62,"WIN 女性领导力"); rule(s,2.1)
bullets(s,0.64,2.5,7.7,[
    [("女性“首席点灯人”主讲，",CREAMTX,False,15),("尤其触达留守女童",WHITE,True,15)],
    "核心信息：“女孩，也可以造芯片、当工程师、当领导者”",
    [("把领导力翻译给孩子：",CREAMTX,False,15),("先点亮自己，再去照亮别人",GOLD_SOFT,True,15)],
    [("安排",CREAMTX,False,15),("女工程师专门结对女学生",WHITE,True,15),("，让榜样可触摸",CREAMTX,False,15)],
    "引用 WIN 理念：多元的头脑汇聚，才有开创性的进步"],size=15,color=CREAMTX,gap=8)
_,tf=box(s,0.64,5.55,7.9,1.2); para(tf,True,[("设计点：",GOLD_SOFT,True,13),("WIN 与“芯火”主题严丝合缝——领导力的本质，就是“点亮人”。女性榜样不再是台上遥远的名字，而是孩子今天真实对话过的人。",CREAMTX,False,13)],lh=1.35)
foot(s,12,True)

# 13 芯·燃 overview
s=slide(); kicker(s,"DAY 2 · 13:40–15:10 · 操场 / 篮球场"); title(s,"芯·燃 ｜ 火，越烧越旺")
nx=tag(s,0.64,1.62,"We Care"); tag(s,nx,1.62,"We Deliver · 道具留校"); rule(s,2.1)
data=[["玩法（工程师与孩子混队）","道具（留校）"],
 ["彩虹伞·点亮星空：齐力把“星星”软球抛上伞面","大彩虹伞、软球"],
 ["芯片运输大挑战：平衡盘端“芯片”过障碍","平衡盘、沙包"],
 ["毛毛虫赛跑：多人同行，谁都不掉队","毛毛虫/跳袋"],
 ["能量传递接力：球拍运”能量球”","羽毛球拍、软球"],
 ["芯火大拔河：最有气氛的一战","安全拔河绳"],
 ["篮球友谊赛 + 定点投篮：人人能得分","篮球 ×6"]]
table(s,0.64,2.4,7.2,data,[4.6,2.0],fs=12,headfs=12.5,rh=0.42)
card(s,8.1,2.4,4.55,3.7,fill=CREAM)
_,tf=box(s,8.35,2.6,4.05,3.4); para(tf,True,"🎁 「ASM 芯火快乐运动角」",size=16,color=PURPLE,bold=True,after=6)
para(tf,False,"所有器材打包捐给学校，活动后建一个长期运动角：",size=13,color=INK,after=6,lh=1.3)
para(tf,False,"彩虹伞×2 · 软球 · 篮球×6 · 足球×4 · 跳绳×20 · 拔河绳×2 · 沙包/平衡盘 · 呼啦圈×20 · 飞盘×20 · 羽毛球 · 收纳袋+标牌",size=12,color=INK,after=8,lh=1.35)
para(tf,False,[("巧思：",GOLDINK,True,12.5),("孩子和工程师一起挂上“运动角“牌——”以后你们每次来玩，就是我们还在陪你们。”",GOLDINK,False,12.5)],lh=1.3)
foot(s,13)

# 13·1 / 13·2 / 13·3 详解
sport=[("①/③",[("sp-umbrella","① 彩虹伞 · 点亮星空","全员抓住大彩虹伞边缘，把软球“星星”抛上伞面，齐心抖动让它跳跃不落；听口令一起高抛，满天星腾空。","👥 全员·40人　⏱ 10′　🎒 彩虹伞·软球","亮点：零门槛·强协作，直扣“星河/点亮”。"),
              ("sp-chip","② 芯片运输大挑战","4-5 人一队接力，用平衡盘端着“易碎芯片”(沙包)走过障碍区，全程不能用手扶、不能掉落；最快且零掉落者胜。","👥 4-5人/队　⏱ 15′　🎒 平衡盘·沙包·桩筒","亮点：呼应 ASM 高精度搬运，玩专注与配合。")]),
       ("②/③",[("sp-caterpillar","③ 毛毛虫赛跑","多人坐进“毛毛虫”道具或一字搭肩蹲走，齐步前进，谁掉队就重来；先到终点的队获胜。","👥 6-8人/队　⏱ 12′　🎒 毛毛虫道具/无道具版","亮点：“一个都不能少”，团队感拉满。"),
              ("sp-relay","④ 能量传递接力","用球拍(或勺)托着“能量球”跑过赛道交给下一棒，球掉了回到掉落点继续；最快队获胜。","👥 5-6人/队　⏱ 12′　🎒 球拍/勺·软球","亮点：简单刺激、人人上手，紧张又好笑。")]),
       ("③/③",[("sp-tug","⑤ 芯火大拔河","工程师+孩子混合两队，中线插“芯火旗”，听哨拔河，把旗拉过己方线者胜；三局两胜。","👥 两队各15-20人　⏱ 12′　🎒 安全拔河绳·中线旗","亮点：最团结的一战，全场欢呼。"),
              ("sp-basket","⑥ 篮球友谊赛 + 定点投篮","混合组队半场友谊赛，规则放宽重在参与；穿插定点投篮，每人3投，进球积分累计到队伍。","👥 两队各5-8人轮换　⏱ 余下时间　🎒 篮球×6","亮点：大人孩子平起平坐当队友。")])]
for idx,(seq,acts) in enumerate(sport):
    s=slide(); kicker(s,seq.join(["DAY 2 · 芯·燃 玩法详解　",""])); title(s,"趣味运动会 · 逐项拆解",size=27); rule(s,1.7)
    for j,(img,t,desc,meta,hi) in enumerate(acts):
        l=0.64+j*6.15
        card(s,l,2.0,5.85,4.6)
        pic_panel(s,l+0.22,2.2,5.4,1.85,img)
        _,tf=box(s,l+0.3,4.18,5.3,2.3)
        para(tf,True,t,size=16,color=PURPLE,bold=True,after=4)
        para(tf,False,desc,size=12.5,color=INK,after=5,lh=1.3)
        para(tf,False,meta,size=11,color=MUTED,after=3)
        para(tf,False,hi,size=11.5,color=GOLDINK,bold=True,lh=1.25)
    foot(s,["13·1","13·2","13·3"][idx])

# 14 芯·山 overview
s=slide(PURPLE_DEEP,"scene-night"); kicker(s,"DAY 2 · 15:10–16:00 · 收尾压轴"); title(s,"芯·山 ｜ 一束芯火，照亮一座山",color=WHITE)
nx=tag(s,0.64,1.62,"We Innovate"); nx=tag(s,nx,1.62,"We Deliver"); tag(s,nx,1.62,"We Care"); rule(s,2.1)
fin=[("① 共绘星河（25′）","预印画布 + 金银马克笔（速干不洒、无需水电）。每对画一颗 logo 星形的“星”+梦想，连成星河。永久留校。"),
     ("② 传光仪式（12′）","人手一盏电池 LED 小灯（无明火）。一束传四十束，齐念 slogan。小灯送给孩子。"),
     ("③ 时间胶囊（13′）","孩子写《给未来的自己》，工程师写给搭档，封存约定回访开启。可选「时光手账」带回家。")]
for i,(h,b) in enumerate(fin):
    l=0.64+i*4.05; card(s,l,2.4,3.8,2.5,fill=PURPLE_MID,ln=PURPLE_MID)
    _,tf=box(s,l+0.22,2.6,3.4,2.1); para(tf,True,h,size=15.5,color=GOLD_SOFT,bold=True,after=6)
    para(tf,False,b,size=12.5,color=CREAMTX,lh=1.3)
_,tf=box(s,0.64,5.15,12,1.0); para(tf,True,[("为什么打动 ASM：",GOLD_SOFT,True,12.5),("传光让每位工程师亲手完成“点亮”，写信让他们真正走进一个孩子，胶囊给所有人一个必须再回来的理由。",CREAMTX,False,12.5)],lh=1.3)
foot(s,14,True)

# 14·1 共绘星河
s=slide(); kicker(s,"DAY 2 · 芯·山 玩法详解 · 第一幕"); title(s,"共绘「星河长卷」（约 25′）"); rule(s,1.78)
pic_panel(s,0.64,2.1,5.6,4.3,"f-canvas",panel=IVORY)
bullets(s,6.5,2.3,6.1,["铺开预印画布（淡星形栅格 + 山的剪影 + slogan）",
    [("每对搭档用金/银马克笔描一颗 ",INK,False,15),("ASM 星形的星",INK,True,15),("，写名字 + 一个梦想词",INK,False,15)],
    [("用线把自己的星和邻座的星连起来 → 全场连成",INK,False,15),("照亮山顶的星河",INK,True,15)]],size=15,gap=7)
_,tf=box(s,6.5,4.25,6.1,0.5); para(tf,True,"🎒 星河画布×1 · 金银马克笔×20 · 星形模板×20　｜　🏫 永久挂校墙",size=12,color=MUTED)
_,tf=box(s,6.5,4.85,6.1,1.0); para(tf,True,[("话术：",PURPLE,True,13),("“你画的不是普通的星——它有你的名字和梦想。看，它和大家的星连在一起，就照亮了整座山。”",INK,False,13)],lh=1.3)
tag(s,6.5,5.95,"We Innovate"); tag(s,8.4,5.95,"We Care",care=True)
foot(s,"14·1")

# 14·2 传光仪式
s=slide(); kicker(s,"DAY 2 · 芯·山 玩法详解 · 第二幕"); title(s,"传光仪式（约 12′）"); rule(s,1.78)
pic_panel(s,0.64,2.1,5.6,4.3,"f-pass")
bullets(s,6.5,2.3,6.1,[[("每人发一盏",INK,False,15),("电池 LED 小灯",INK,True,15),("（无明火、无需供电）",INK,False,15)],
    [("由女性“首席点灯人”手中",INK,False,15),("唯一亮着的灯",INK,True,15),("开始，手手相传",INK,False,15)],
    [("直到 40 人全部亮起，举灯齐念 ",INK,False,15),("“一束芯火，照亮一座山”",INK,True,15)]],size=15,gap=7)
_,tf=box(s,6.5,4.3,6.1,0.5); para(tf,True,"🎒 电池 LED 小灯×40　｜　🏫 留存：小灯全部送孩子带回家",size=12,color=MUTED)
_,tf=box(s,6.5,4.9,6.1,0.9); para(tf,True,[("话术：",PURPLE,True,13),("“光不会因为分给别人而变少——它只会让更多人亮起来。这盏灯送给你，它会一直亮着。”",INK,False,13)],lh=1.3)
_,tf=box(s,6.5,5.85,6.1,0.4); para(tf,True,"已砍掉 UV 暗房依赖，简易学校可稳稳跑完。",size=12,color=GOLDINK,bold=True)
foot(s,"14·2")

# 14·3 胶囊+手账
s=slide(); kicker(s,"DAY 2 · 芯·山 玩法详解 · 第三幕"); title(s,"时间胶囊 + 时光手账（约 13′）",size=27); rule(s,1.7)
pic_panel(s,0.64,2.0,12.05,1.5,"f-capsule")
card(s,0.64,3.65,5.9,3.0)
_,tf=box(s,0.9,3.85,5.4,2.7); para(tf,True,"🕰 时间胶囊（集体记忆 · 留校）",size=16,color=PURPLE,bold=True,after=6)
bullets(s,0.9,4.35,5.4,["孩子写《给未来的自己》","每位工程师写一封给搭档孩子","一起封进胶囊盒，交校方保管，约定下次回访开启"],size=13.5,gap=5)
card(s,6.75,3.65,5.9,3.0,fill=CREAM)
_,tf=box(s,7.0,3.85,5.4,2.7); para(tf,True,"📖 时光手账（个人陪伴 · 带回家）",size=16,color=PURPLE,bold=True,after=4)
para(tf,False,"每个孩子获赠一本「芯火时光手账」（深紫×暖金、印 ASM 星形母题）：",size=12.5,color=INK,after=4,lh=1.3)
bullets(s,7.0,4.85,5.4,["内页：今日打卡、给未来的自己、梦想清单、双师课堂记录页、成长贴纸",
    "工程师在搭档手账上亲笔留一页寄语","带回家持续书写，衔接后续线上双师课堂"],size=12.5,gap=4)
_,tf=box(s,0.64,6.75,9,0.4); para(tf,True,[("区别：",PURPLE,True,12.5),("胶囊＝集体记忆留校；手账＝个人陪伴带回家——二者并存。",INK,False,12.5)])
foot(s,"14·3")

# 15 推荐方案2
s=slide(); kicker(s,"推荐方案 2（可选预案）"); title(s,"两座城，一片天 ｜ 上海孩子同行"); rule(s,1.78)
_,tf=box(s,0.64,2.0,6,0.4); para(tf,True,"核心设计",size=18,color=PURPLE,bold=True,after=2)
bullets(s,0.64,2.5,6.0,["不是城里孩子来“看”山里孩子，而是两个世界互相照亮",
    "三人组结对：1 上海孩子 + 1 木贡孩子 + 1 工程师",
    "角色反转：木贡孩子当东道主/向导，在带路·自然·运动·篮球中是“专家”，在主场收获自信",
    "收尾升华：两座城的星，连进同一条星河"],size=14,gap=6)
_,tf=box(s,6.95,2.0,6,0.4); para(tf,True,"价值 & 注意",size=18,color=PURPLE,bold=True,after=2)
bullets(s,6.95,2.5,5.9,["✅ 故事性、传播性、ESG 叙事更强（消除城乡认知鸿沟）",
    "✅ 双向成长：上海孩子收获感恩与坚韧",
    "⚠ 增加招募、监护、保险与安全成本",
    "⚠ 做好引导，规避任何“贫富对比/猎奇”观感（已在角色设计上规避）"],size=14,gap=6)
foot(s,15)

# 16 人员分工
s=slide(); kicker(s,"执行 · 人员分工　ASM 约 20 人"); title(s,"每个人既是“搭档”，也有一份岗位"); rule(s,1.78)
data=[["角色","人数","职责"],
 ["总协调 / 主持","1","全程节奏、串场、对接校方与时间把控"],
 ["首席点灯人（女性工程师）","1","Day1 点火 → 芯·她 WIN 主讲 → 收尾长期之约"],
 ["STEM 讲师","2–3","科研讲座主讲 + 实验带教总指挥"],
 ["运动会裁判 / 安全员","2","趣味赛与篮球赛组织、规则、安全"],
 ["影像（摄影 + 摄像）","2","高光时刻记录，供后续传播与 ESG 报告"],
 ["后勤 / 物料","2","物料分发、场地布置、捐赠物资交接"],
 ["1:1 芯火搭档","全体 20","以上岗位叠加在 20 人身上；每人全程陪伴 1 个孩子"]]
table(s,0.64,2.05,12.05,data,[2.6,1.1,6.0],fs=12.5,headfs=13,rh=0.46)
_,tf=box(s,0.64,6.35,12,0.4); para(tf,True,"20 名 ASM 人员 = 20 名搭档，岗位角色叠加分配；确保 20:20 严格 1:1，无孩子掉队。",size=12,color=MUTED)
foot(s,16)

# 17 物料清单
s=slide(); kicker(s,"执行 · 物料采购总清单"); title(s,"物料一览（含数量 · 是否留校）"); rule(s,1.78)
d1=[["仪式 / 课程物料","数量","留校"],["发光拼合徽章","40","—"],["STEM 点灯实验套件","20","✓"],
 ["芯火树灯组","1","—"],["电池 LED 小灯","40","✓"],["星河画布+金银马克笔+模板","1+20","✓"],
 ["时间胶囊盒+卡片/时光手账","1+40","✓"],["助学金 + 科教物资","—","✓"]]
d2=[["运动 / 通用物料","数量","留校"],["彩虹伞","2","✓"],["篮球 / 足球","6 / 4","✓"],
 ["跳绳/呼啦圈/飞盘","20×3","✓"],["拔河绳/沙包/平衡盘","2/批","✓"],["横幅·桌签·物料牌","—","部分"],
 ["影像设备·备用电源","—","—"],["医药急救箱","2","—"]]
table(s,0.64,2.05,6.0,d1,[3.6,1.0,0.8],fs=11.5,headfs=12,rh=0.46)
table(s,6.95,2.05,5.7,d2,[3.4,1.0,0.8],fs=11.5,headfs=12,rh=0.46)
_,tf=box(s,0.64,6.45,12,0.4); para(tf,True,"原则：凡能留下的，全部留给学校——让陪伴在我们离开后继续。",size=12,color=MUTED)
foot(s,17)

# 18 动线+应急
s=slide(); kicker(s,"执行 · 场地动线 & 应急预案"); title(s,"把“意外”提前想到"); rule(s,1.78)
card(s,0.64,2.1,5.9,3.6)
_,tf=box(s,0.9,2.3,5.4,3.2); para(tf,True,"场地动线",size=17,color=PURPLE,bold=True,after=6)
bullets(s,0.9,2.8,5.4,["教室：开幕/捐赠/STEM/WIN/时间胶囊","操场：趣味运动会/传光仪式/共绘星河",
    "篮球场：篮球友谊赛 + 定点投篮","动线就近，午餐在食堂，全程 1:1 不分散"],size=14,gap=6)
card(s,6.75,2.1,5.9,3.6)
_,tf=box(s,7.0,2.3,5.4,3.2); para(tf,True,"应急预案",size=17,color=PURPLE,bold=True,after=6)
bullets(s,7.0,2.8,5.4,["天气：雨天运动会改室内；星河、传光移入教室","电力：LED 全电池供电、已砍 UV，不依赖现场供电",
    "医疗：2 个急救箱 + 预确认最近医院与路线","情绪：班主任全程陪同，留守儿童情绪引导","人数浮动：1:1 弹性分组，多备 2 套物料"],size=14,gap=5)
_,tf=box(s,0.64,5.9,12,0.5); para(tf,True,[("落地铁律：",PURPLE,True,13),("所有环节”不依赖暗房 / 不依赖稳定电力 / 不依赖网络“，简易学校也能稳稳跑完。",INK,False,13)])
foot(s,18)

# 19 儿童保护合规
s=slide(); kicker(s,"执行 · 儿童保护与合规（不可省略）"); title(s,"把孩子的安全，放在创意之前"); rule(s,1.78)
card(s,0.64,2.1,5.9,3.0)
_,tf=box(s,0.9,2.3,5.4,2.6); para(tf,True,"儿童保护",size=17,color=PURPLE,bold=True,after=6)
bullets(s,0.9,2.8,5.4,["肖像 / 影像须监护人书面授权，对外传播去隐私化","志愿者签儿童保护行为准则：不与儿童单独相处","全程班主任 / 校方在场陪同"],size=14,gap=6)
card(s,6.75,2.1,5.9,3.0)
_,tf=box(s,7.0,2.3,5.4,2.6); para(tf,True,"安全与物料合规",size=17,color=PURPLE,bold=True,after=6)
bullets(s,7.0,2.8,5.4,["发光物料一律电子烛火、无明火，3C/无毒认证、低压安全","定制件预留打样周期，备足备用件 + 故障兜底","2 个急救箱 + 预确认最近医院路线；保险覆盖全员"],size=14,gap=6)
card(s,0.64,5.3,12.05,1.15,fill=CREAM)
_,tf=box(s,0.9,5.5,11.5,0.85); para(tf,True,[("待官方确认：",GOLDINK,True,13),("本方案中 ASM 价值观 / 使命金句 / WIN 理念 / Logo 描述 / 帮扶年限等，系依据简报整理，正式定稿前请以 ASM 品牌 · CSR 官方表述为准。",INK,False,13)],lh=1.3)
foot(s,19)

# 20 倒排期
s=slide(PURPLE_DEEP); kicker(s,"执行 · 筹备倒排期 Checklist"); title(s,"从今天到现场，倒着排",color=WHITE); rule(s,1.78)
data=[["节点","关键动作"],
 ["T-30 天","确认行程/人数；校方对接（学生名单、场地、电力）；STEM 课程包与讲师定稿；物料采购下单"],
 ["T-14 天","定制件到货；讲师试讲彩排；影像与传播脚本；女性点灯人内容定稿"],
 ["T-7 天","全流程桌面推演（含应急）；物料清点装箱；1:1 搭档名单与岗位分工锁定；安全/保险确认"],
 ["T-1 天","抵达布场预演（芯火树、星河画布、运动场地）；设备充电；Day1 破冰最终走位"],
 ["现场","按 Run-of-show 执行；总协调控时；影像抓高光；当晚团队复盘"]]
table(s,0.64,2.1,12.05,data,[1.3,8.0],fs=12.5,headfs=13,rh=0.62)
foot(s,20,True)

# 21 影响+回响
s=slide(); kicker(s,"影响 · 认知拔升 & 长期回响"); title(s,"这一天，会留在孩子的认知里很久"); rule(s,1.78)
data=[["升级前","经由","升级后"],["我是被资助的孩子","带路·回赠","我是主人，也能给予"],
 ["高科技离我很远","点亮 LED","我也是工程师"],["女孩读书没用","芯·她 WIN","女孩也能当领导"],
 ["我没什么潜力","芯火隐喻","我是一束有光的芯"],["帮一次就走","时间胶囊","有一个长期的约定"]]
table(s,0.64,2.05,6.4,data,[2.0,1.4,2.2],fs=12,headfs=12.5,rh=0.5)
card(s,7.4,2.05,5.25,1.6)
_,tf=box(s,7.65,2.25,4.8,1.3); para(tf,True,"长期回响机制",size=15,color=PURPLE,bold=True,after=4)
para(tf,False,"时间胶囊回访开启 · 线上双师课堂延续 · 择优上海科创夏令营 · 星河墙与运动角永久留校",size=12.5,color=INK,lh=1.3)
card(s,7.4,3.85,5.25,1.6,fill=CREAM)
_,tf=box(s,7.65,4.05,4.8,1.3); para(tf,True,"成效追踪（KPI）",size=15,color=PURPLE,bold=True,after=4)
para(tf,False,"结对覆盖 20/20 · 人均 1 件科创作品 · 40 封“未来之信” · 回访开启率 · 影像与 ESG 传播产出",size=12.5,color=INK,lh=1.3)
_,tf=box(s,0.64,6.3,12,0.5); para(tf,True,"从“给物资“升级到”给视野、给榜样、给自我效能、给成长型思维”——这比任何物资都留得久。",size=12.5,color=MUTED)
foot(s,21)

# 22 CLOSING
s=slide(PURPLE_DEEP,"scene-night")
kicker(s,"ASM · WE INNOVATE · WE DELIVER · WE CARE",top=1.4)
_,tf=box(s,0.9,2.2,11.5,3.0)
para(tf,True,"ASM 不是去“帮”这些孩子，",size=32,color=WHITE,bold=True,after=2,lh=1.25)
para(tf,False,"而是去做它最擅长的事——",size=32,color=WHITE,bold=True,after=2,lh=1.25)
para(tf,False,[("把一颗颗被埋没的“芯”，连接、点亮，然后让它发光。",GOLD_SOFT,True,32)],after=0,lh=1.25)
_,tf=box(s,0.9,5.5,11,0.8); para(tf,True,"一束芯火，照亮一座山。",size=24,color=GOLD_SOFT,bold=True)
foot(s,22,True)

prs.save("ASM_芯火计划_执行方案.pptx")
print("SAVED", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
PY = None
